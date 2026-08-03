import asyncio
import logging
import os
import re
import traceback
from datetime import datetime, timedelta, timezone
from decimal import Decimal
from io import BytesIO

import ai
import boto3
import groups
import templates
from cache import CA_LINKS, CHORDS, MP3, PIANO, SCORES, SGM_LINKS, SONGS, TITLES, VIDEOS
from lookup import SONGS_LOOKUP, TITLES_LOOKUP
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from rapidfuzz import fuzz, process
from telegram import (
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    InputMediaAudio,
    InputMediaPhoto,
    KeyboardButton,
    ReplyKeyboardMarkup,
    ReplyKeyboardRemove,
    Update,
    constants,
)
from telegram.error import TelegramError
from telegram.ext import (
    Application,
    CallbackQueryHandler,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    filters,
)
from unidecode import unidecode

logger = logging.getLogger()
logger.setLevel("INFO")

BOT_TOKEN = os.getenv("BOT_TOKEN")
VIDEOS_S3_BUCKET = os.getenv("VIDEOS_S3_BUCKET")

app = Application.builder().token(BOT_TOKEN).build()
dynamodb = boto3.resource("dynamodb")


def saveLog(user, event, request, response):
    now = datetime.now(timezone.utc)
    dynamodb.Table("tsms_logs").put_item(
        Item={
            "user_id": user.id,
            "name": user.full_name,
            "username": user.username,
            "timestamp": Decimal(str(now.timestamp())),
            "timestamp_iso": now.astimezone(timezone(timedelta(hours=8))).isoformat(),
            "event": event,
            "request": request,
            "response": response,
        }
    )


def getDbUser(user):
    dbUser = dynamodb.Table("tsms_users").get_item(Key={"id": user.id})
    return dbUser


async def require_registration(update: Update, dbUser, action) -> bool:
    if "Item" in dbUser:
        return True
    await update.effective_chat.send_message("Press /start to begin")
    saveLog(update.effective_user, "NOT_REGISTERED", action, None)
    return False


async def reply_and_log(update: Update, text, event, request=None, response=None, **kwargs):
    await update.message.reply_html(text, **kwargs)
    saveLog(update.effective_user, event, request, response)


async def send_and_log(update: Update, text, event, request=None, response=None, **kwargs):
    await update.effective_chat.send_message(
        text=text, parse_mode=constants.ParseMode.HTML, **kwargs
    )
    saveLog(update.effective_user, event, request, response)


async def updateState(update: Update, dbUser) -> None:
    user = update.effective_user
    phone = dbUser["Item"]["phone"].lstrip("+")
    state = dbUser["Item"]["state"]
    if state != templates.current_version:
        dynamodb.Table("tsms_users").update_item(
            Key={"id": user.id},
            UpdateExpression="SET phone = :phone, #st = :state",
            ExpressionAttributeValues={
                ":phone": phone,
                ":state": templates.current_version,
            },
            ExpressionAttributeNames={"#st": "state"},
        )
        await reply_and_log(update, templates.changelog, "CHANGELOG", response=templates.current_version)


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    contact_keyboard = KeyboardButton(text="REGISTER", request_contact=True)
    reply_markup = ReplyKeyboardMarkup(
        [[contact_keyboard]],
        one_time_keyboard=True,
        input_field_placeholder="Use button below to register",
    )
    await reply_and_log(update, templates.start, "START", reply_markup=reply_markup)


async def contact(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user = update.effective_user
    contact = update.message.contact
    phone = contact.phone_number.lstrip("+")
    if user.id == contact.user_id and phone.startswith(templates.allowed_phone):
        dynamodb.Table("tsms_users").put_item(
            Item={
                "id": user.id,
                "name": user.full_name,
                "phone": phone,
                "state": templates.current_version,
            }
        )
        await reply_and_log(
            update, templates.welcome, "CONTACT", request=phone, response="Allowed",
            reply_markup=ReplyKeyboardRemove(),
        )
    else:
        await reply_and_log(
            update,
            "Sorry, our automated security check does not allow you to use this bot.",
            "CONTACT", request=phone, response="Blocked",
            reply_markup=ReplyKeyboardRemove(),
        )


async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    dbUser = getDbUser(update.effective_user)
    if not await require_registration(update, dbUser, "HELP"):
        return
    await reply_and_log(
        update, templates.examples, "HELP",
        reply_markup=ReplyKeyboardRemove(), protect_content=True,
    )


def make_button(song_number, option, callback_prefix, button_text):
    if option:
        return [
            [
                InlineKeyboardButton(
                    button_text, callback_data=f"{callback_prefix} {song_number}"
                )
            ]
        ]
    return []


async def send_link_buttons(update: Update, links, source, song_number) -> None:
    keyboard = [
        [InlineKeyboardButton(key, url=value)] for key, value in links.items()
    ]
    await send_and_log(
        update, f"Links provided by {source}", "LINKS", request=source, response=song_number,
        disable_web_page_preview=True,
        reply_markup=InlineKeyboardMarkup(keyboard),
    )


async def send_song(
    update: Update, context: ContextTypes.DEFAULT_TYPE, song_number, dbUser, event, request
) -> None:
    active_group_id = groups.get_active_group(update.effective_user.id, dbUser)
    keyboard = []
    keyboard.extend(
        make_button(song_number, song_number in CHORDS, "CHORDS", "🎸 Guitar Chords")
    )
    keyboard.extend(
        make_button(song_number, song_number in SCORES, "SCORE", "🎼 Piano Score")
    )
    keyboard.extend(
        make_button(song_number, song_number in MP3, "MP3", "🔊 MIDI Soundtrack")
    )
    keyboard.extend(
        make_button(
            TITLES[song_number],
            TITLES[song_number] in PIANO,
            "PIANO",
            "🎹 Piano Recording (Wilds)",
        )
    )
    keyboard.extend(
        make_button(
            TITLES[song_number],
            TITLES[song_number] in VIDEOS,
            "VIDEO",
            "🎤  Choir Recording (Lyric Video)",
        )
    )
    lyrics = SONGS.get(song_number)
    if lyrics.count("\n\n") > 0:
        keyboard.extend(make_button(song_number, True, "PPT", "💻 Generate PowerPoint"))
        keyboard.extend(make_button(song_number, True, "EXPLAIN", "💭 Explain Song"))
    keyboard.extend(
        make_button(
            song_number, active_group_id is not None, "GROUP_SEND", "📢 Send to Community"
        )
    )
    await send_and_log(
        update, lyrics, event, request=request, response=f"{song_number} {TITLES[song_number]}",
        disable_web_page_preview=True,
        reply_markup=InlineKeyboardMarkup(keyboard),
    )
    if song_number in CA_LINKS:
        await send_link_buttons(update, CA_LINKS[song_number], "cityalight.com", song_number)
    if song_number in SGM_LINKS:
        await send_link_buttons(update, SGM_LINKS[song_number], "sovereigngracemusic.com", song_number)
    await groups.process_search_event(context, saveLog, update.effective_user, song_number, active_group_id)


async def search(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user = update.effective_user
    dbUser = getDbUser(user)
    if not await require_registration(update, dbUser, "SEARCH"):
        return
    raw_message = update.message.text
    message = unidecode(raw_message).replace("\n", " ").strip().upper()
    hint = message.startswith("TSMS")
    if message.isnumeric():  # handle default book
        message = int(message)
        message = "TSMS " + str(message)
    song_number = None
    results = []
    if message in SONGS:  # match song number
        song_number = message
    else:  # match title
        alpha_only = re.compile("[^A-Z ]")
        clean_message = alpha_only.sub("", message).strip()
        if clean_message in TITLES_LOOKUP:
            results = TITLES_LOOKUP.get(clean_message).copy()
            song_number = results.pop(0)
        else:  # search
            if len(clean_message) > 200:
                await reply_and_log(
                    update, "<i>Please shorten your search</i>", "SEARCH_TOO_LONG", request=raw_message
                )
                return
            await update.message.reply_chat_action(constants.ChatAction.TYPING)
            query = process.extract(
                clean_message,
                SONGS_LOOKUP,
                scorer=fuzz.partial_ratio,
                score_cutoff=85,
                limit=10,
            )
            results = [t[2] for t in query]

    if song_number:
        await send_song(update, context, song_number, dbUser, "SEARCH_HIT", raw_message)
        if hint:
            await reply_and_log(
                update,
                "<i>Hint: Get to TSMS songs faster by typing the song number without 'TSMS'</i>",
                "SEARCH_HINT", request=raw_message, response=song_number,
            )
    if results:
        keyboard = []
        for number in results:
            keyboard.extend(
                make_button(number, True, "SONG", f"{number} {TITLES[number]}")
            )
        if song_number:
            await reply_and_log(
                update, "<i>This song is also found in:</i>", "SEARCH_ALSO_FOUND",
                request=raw_message, response=f"{len(results)} results",
                reply_markup=InlineKeyboardMarkup(keyboard),
            )
        else:
            await reply_and_log(
                update, "<i>Showing up to 10 search results:</i>", "SEARCH_RESULTS",
                request=raw_message, response=f"{len(results)} results",
                reply_markup=InlineKeyboardMarkup(keyboard),
            )
    elif not song_number:
        await reply_and_log(
            update, "<i>No matches found</i>\n\nType /help for instructions", "SEARCH_NONE",
            request=raw_message,
        )
    await updateState(update, dbUser)


def make_ppt(song_number):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    title = TITLES.get(song_number)
    text = SONGS.get(song_number)
    text = text.split("\n\n")
    text.pop(0)
    text = list(filter(None, text))

    originallen = len(text)
    chorus = None
    for i in range(originallen):
        stanza = text[i]
        if stanza.startswith("Chorus:") or stanza.startswith("Refrain:"):
            chorus = i
            break
    if chorus is not None:
        i = chorus + 2
        while True:
            text.insert(i, stanza)
            i += 2
            if i > len(text):
                break

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    txBox = slide.shapes.add_textbox(0, 0, Inches(16), Inches(9))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.add_paragraph()
    p.text = title + "\n(" + song_number + ")"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    for i in range(len(text)):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        txBox = slide.shapes.add_textbox(Inches(15), 0, Inches(1), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "{}/{}".format(i + 1, len(text))
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(255, 255, 255)

        txBox = slide.shapes.add_textbox(0, 0, Inches(16), Inches(9))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = tf.add_paragraph()
        p.text = text[i].strip()
        p.font.size = Pt(48)
        if song_number.startswith("C "):
            p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    pptxfile = BytesIO()
    pptxfile.name = song_number + ".pptx"
    prs.save(pptxfile)
    pptxfile.seek(0)
    return pptxfile


async def answer_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user = update.effective_user
    dbUser = getDbUser(user)
    if not await require_registration(update, dbUser, "CALLBACK"):
        return
    query = update.callback_query
    data = query.data
    if data.startswith("SONG "):
        song_number = data.replace("SONG ", "")
        await send_song(update, context, song_number, dbUser, "CALLBACK", "SONG")
    elif data.startswith("CHORDS "):
        song_number = data.replace("CHORDS ", "")
        await send_and_log(
            update, f"<pre>{CHORDS[song_number]}</pre>", "CALLBACK", request="CHORDS", response=song_number
        )
    elif data.startswith("SCORE "):
        song_number = data.replace("SCORE ", "")
        await update.effective_chat.send_action(constants.ChatAction.UPLOAD_PHOTO)
        saveLog(user, "CALLBACK", "SCORE", song_number)
        media_list = SCORES[song_number]
        media = [InputMediaPhoto(media=p) for p in media_list]
        await update.effective_chat.send_media_group(media=media)
    elif data.startswith("MP3 "):
        song_number = data.replace("MP3 ", "")
        await update.effective_chat.send_action(constants.ChatAction.UPLOAD_DOCUMENT)
        saveLog(user, "CALLBACK", "MP3", song_number)
        media_list = MP3[song_number]
        media = [InputMediaAudio(media=a) for a in media_list]
        await update.effective_chat.send_media_group(media=media)
    elif data.startswith("PIANO "):
        song_title = data.replace("PIANO ", "")
        await update.effective_chat.send_action(constants.ChatAction.UPLOAD_DOCUMENT)
        saveLog(user, "CALLBACK", "PIANO", song_title)
        reference = PIANO[song_title]
        await update.effective_chat.send_audio(audio=reference, protect_content=True)
    elif data.startswith("VIDEO "):
        song_title = data.replace("VIDEO ", "")
        await update.effective_chat.send_action(constants.ChatAction.UPLOAD_VIDEO)
        saveLog(user, "CALLBACK", "VIDEO", song_title)
        objects = VIDEOS[song_title]
        total = len(objects)
        await query.answer(text=f"Loading {str(total)} video(s), please wait...")
        count = 1
        for obj in objects:
            video = VIDEOS_S3_BUCKET + obj
            caption = (
                f"{song_title.title()} {count}" if total > 1 else song_title.title()
            )
            await update.effective_chat.send_video(video=video, caption=caption)
            count += 1
    elif data.startswith("PPT "):
        song_number = data.replace("PPT ", "")
        await update.effective_chat.send_action(constants.ChatAction.UPLOAD_DOCUMENT)
        saveLog(user, "CALLBACK", "PPT", song_number)
        ppt = make_ppt(song_number)
        await update.effective_chat.send_document(document=ppt)
    elif data.startswith("EXPLAIN "):
        song_number = data.replace("EXPLAIN ", "")
        await update.effective_chat.send_action(constants.ChatAction.TYPING)
        saveLog(user, "CALLBACK", "EXPLAIN", song_number)
        await query.answer(text="Thinking...")
        response = ai.explainSong(SONGS.get(song_number))
        await update.effective_chat.send_message(
            response, parse_mode=constants.ParseMode.HTML
        )
    elif data.startswith("GROUP_SEND "):
        song_number = data.replace("GROUP_SEND ", "")
        saveLog(user, "CALLBACK", "GROUP_SEND", song_number)
        active_group_id = groups.get_active_group(user.id, dbUser)
        if active_group_id is None:
            await query.answer(text="You are currently not in a Community", show_alert=True)
            saveLog(user, "GROUP_SEND_BLOCKED", song_number, "NOT_IN_COMMUNITY")
            return
        members = groups.try_mark_sent(active_group_id, song_number)
        if members is None:
            await query.answer(text="Song has already been sent.", show_alert=True)
            saveLog(user, "GROUP_SEND_BLOCKED", song_number, "ALREADY_SENT")
            return
        lyrics = SONGS.get(song_number) + templates.group_sent_caption.format(
            name=user.full_name
        )

        async def send_to_member(member_id):
            log_user = groups.get_log_user(member_id)
            try:
                await context.bot.send_message(
                    chat_id=member_id,
                    text=lyrics,
                    parse_mode=constants.ParseMode.HTML,
                    disable_web_page_preview=True,
                )
            except TelegramError:
                logger.exception("Failed to deliver GROUP_SEND to %s", member_id)
                saveLog(log_user, "GROUP", active_group_id, f"DELIVERY_FAILED {song_number}")
                return
            saveLog(log_user, "GROUP", active_group_id, f"RECEIVED {song_number}")

        await asyncio.gather(*(send_to_member(member_id) for member_id in members))
        saveLog(user, "GROUP", active_group_id, f"SEND {song_number}")
    else:
        await query.answer(text="This feature is not available")
        saveLog(user, "CALLBACK_INVALID", data, None)
        return
    await query.answer()


async def leave(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user = update.effective_user
    dbUser = getDbUser(user)
    if not await require_registration(update, dbUser, "LEAVE"):
        return
    active_group_id = groups.get_active_group(user.id, dbUser)
    if active_group_id is None:
        await reply_and_log(update, templates.group_not_in_community, "LEAVE_BLOCKED")
        return
    groups.leave_group(user.id, active_group_id)
    await reply_and_log(update, templates.group_left, "GROUP", request=active_group_id, response="LEAVE")


async def tg_bot_main(bot_app, event):
    async with bot_app:
        logger.info("PROCESSING UPDATE: %s", event)
        await bot_app.process_update(Update.de_json(event, bot_app.bot))


def lambda_handler(event, context):
    if "healthCheck" in event:
        return {"statusCode": 200}
    try:
        asyncio.run(tg_bot_main(app, event))
    except Exception as e:
        traceback.print_exc()
        return {"statusCode": 500}
    return {"statusCode": 200}


app.add_handler(CommandHandler("start", start))
app.add_handler(MessageHandler(filters.CONTACT, contact))
app.add_handler(CommandHandler("help", help_command))
app.add_handler(CommandHandler("leave", leave))
app.add_handler(MessageHandler(filters.TEXT, search))
app.add_handler(CallbackQueryHandler(answer_callback))
